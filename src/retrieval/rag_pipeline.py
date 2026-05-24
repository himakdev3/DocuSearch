"""
Production-Grade RAG Pipeline with Citation Logic

This module implements a complete RAG system with:
- PDF document processing and chunking
- Semantic embeddings generation
- Vector similarity search
- LLM-based answer generation with citations
- Comprehensive exception handling
"""

import re
import logging
from pathlib import Path
from typing import List, Dict, Tuple, Optional, Any
from dataclasses import dataclass, field
from datetime import datetime

import numpy as np
import PyPDF2
import docx
import faiss
from sentence_transformers import SentenceTransformer
from tqdm import tqdm

# Import custom exceptions
try:
    from src.utils.exceptions import (
        ModelLoadError, DocumentProcessingError,
        EmbeddingGenerationError, IndexBuildError, SearchError
    )
except ImportError:
    # Fallback if exceptions module not available
    class ModelLoadError(Exception):
        pass

    class DocumentProcessingError(Exception):
        pass

    class EmbeddingGenerationError(Exception):
        pass

    class IndexBuildError(Exception):
        pass

    class SearchError(Exception):
        pass

logger = logging.getLogger(__name__)

_QUERY_STOPWORDS = {
    "a", "an", "and", "are", "as", "at", "be", "by", "for", "from", "how",
    "i", "in", "is", "it", "me", "of", "on", "or", "that", "the", "to",
    "was", "what", "when", "where", "which", "who", "why", "with", "you",
}

_MIN_RESULT_CHARS = 50
_MIN_RESULT_WORDS = 8


@dataclass
class DocumentChunk:
    """Represents a chunk of text from a document with metadata."""
    text: str
    doc_name: str
    page_num: int
    chunk_id: int
    embedding: Optional[np.ndarray] = None
    metadata: Dict[str, Any] = field(default_factory=dict)


class RAGPipeline:
    """
    Production-ready RAG pipeline with advanced features.
    
    Features:
    - Intelligent document chunking with overlap
    - Semantic embeddings using SentenceTransformers
    - FAISS vector index for efficient similarity search
    - Citation tracking with source attribution
    - Comprehensive error handling and validation
    """
    
    def __init__(
        self,
        embedding_model_name: str = "all-MiniLM-L6-v2",
        chunk_size: int = 500,
        chunk_overlap: int = 50,
        use_gpu: bool = False
    ):
        """
        Initialize the RAG pipeline.
        
        Args:
            embedding_model_name: Name of the sentence-transformers model
            chunk_size: Maximum characters per chunk
            chunk_overlap: Characters to overlap between chunks
            use_gpu: Whether to use GPU for embeddings
            
        Raises:
            ModelLoadError: If embedding model fails to load
            ValueError: If invalid parameters provided
        """
        try:
            logger.info(f"Initializing RAG Pipeline with model: {embedding_model_name}")
            
            # Validate parameters
            if chunk_size <= 0:
                raise ValueError("chunk_size must be positive")
            if chunk_overlap < 0 or chunk_overlap >= chunk_size:
                raise ValueError("chunk_overlap must be between 0 and chunk_size")
            
            # Model configuration
            self.chunk_size = chunk_size
            self.chunk_overlap = chunk_overlap
            
            # Initialize embedding model
            device = "cuda" if use_gpu else "cpu"
            self.embedding_model = SentenceTransformer(embedding_model_name, device=device)
            self.embedding_dim = self.embedding_model.get_sentence_embedding_dimension()
            
            # Document storage
            self.documents: List[DocumentChunk] = []
            self.index: Optional[faiss.IndexFlatIP] = None  # Inner product for cosine similarity
            self.embeddings_matrix: Optional[np.ndarray] = None
            self.is_indexed = False
            
            logger.info(f"Pipeline initialized. Embedding dim: {self.embedding_dim}")
            
        except Exception as e:
            logger.error(f"Failed to initialize RAG pipeline: {e}")
            raise ModelLoadError(f"Could not initialize RAG pipeline: {e}") from e
    
    def load_documents_from_files(self, uploaded_files: List) -> int:
        """
        Load and process supported documents from uploaded files.
        
        Args:
            uploaded_files: List of uploaded file objects (Streamlit UploadedFile)
            
        Returns:
            Number of chunks created
            
        Raises:
            DocumentProcessingError: If no documents could be processed
        """
        logger.info(f"Loading {len(uploaded_files)} files")
        
        if not uploaded_files:
            raise DocumentProcessingError("No files provided")
        
        self.documents.clear()
        self.embeddings_matrix = None
        self.index = None
        self.is_indexed = False
        chunk_counter = 0
        failed_files = []
        
        for uploaded_file in tqdm(uploaded_files, desc="Processing documents"):
            try:
                doc_name = uploaded_file.name
                ext = Path(doc_name).suffix.lower()

                if ext == ".pdf":
                    chunk_counter = self._process_pdf_file(uploaded_file, doc_name, chunk_counter, failed_files)
                elif ext in (".docx", ".doc"):
                    chunk_counter = self._process_docx_file(uploaded_file, doc_name, chunk_counter, failed_files)
                elif ext == ".txt":
                    chunk_counter = self._process_txt_file(uploaded_file, doc_name, chunk_counter, failed_files)
                elif ext == ".pptx":
                    chunk_counter = self._process_pptx_file(uploaded_file, doc_name, chunk_counter, failed_files)
                else:
                    logger.warning(f"Unsupported file type: {doc_name}")
                    failed_files.append((doc_name, "Unsupported file type"))
                    continue

            except Exception as e:
                error_msg = str(e)
                logger.error(f"Error processing {uploaded_file.name}: {error_msg}")
                failed_files.append((uploaded_file.name, error_msg))
                continue
        
        # Report failed files
        if failed_files:
            logger.warning(f"Failed to process {len(failed_files)} files: {[f[0] for f in failed_files]}")
        
        if len(self.documents) == 0:
            raise DocumentProcessingError(
                f"No documents could be processed. Failed files: {failed_files}"
            )
        
        logger.info(f"Created {len(self.documents)} chunks from uploaded files")
        return len(self.documents)

    def _append_chunks_from_text(
        self,
        text: str,
        doc_name: str,
        page_num: int,
        chunk_counter: int,
        metadata: Optional[Dict[str, Any]] = None,
    ) -> int:
        """Create chunks from text and append them to document storage."""
        if not text or not text.strip():
            return chunk_counter

        chunk_metadata = dict(metadata or {})
        for chunk_text in self._create_chunks(text):
            self.documents.append(DocumentChunk(
                text=chunk_text,
                doc_name=doc_name,
                page_num=page_num,
                chunk_id=chunk_counter,
                metadata=dict(chunk_metadata),
            ))
            chunk_counter += 1

        return chunk_counter

    def _process_pdf_file(self, pdf_file, doc_name: str, chunk_counter: int, failed_files: list) -> int:
        """Extract text (and OCR images) from a PDF uploaded file and append chunks."""
        try:
            # Read all bytes once so both PyPDF2 and PyMuPDF can use them
            import io
            pdf_bytes = pdf_file.read()
            pdf_reader = PyPDF2.PdfReader(io.BytesIO(pdf_bytes))
            page_previews = self._build_pdf_page_previews(pdf_bytes)
            page_labels = self._extract_pdf_page_labels(pdf_bytes)

            if len(pdf_reader.pages) == 0:
                logger.warning(f"PDF '{doc_name}' has no pages")
                failed_files.append((doc_name, "No pages found"))
                return chunk_counter

            logger.info(f"Processing PDF '{doc_name}' ({len(pdf_reader.pages)} pages)")

            extracted_page_texts: Dict[int, str] = {}
            printed_candidates: Dict[int, int] = {}

            for page_num, page in enumerate(pdf_reader.pages, start=1):
                try:
                    text = page.extract_text()
                    extracted_page_texts[page_num] = text or ""

                    candidate = self._extract_printed_page_candidate(text)
                    if candidate is not None:
                        printed_candidates[page_num] = candidate
                except Exception as e:
                    logger.error(f"Error processing page {page_num} of '{doc_name}': {e}")
                    continue

            printed_page_map = self._resolve_printed_page_map(printed_candidates, len(pdf_reader.pages))

            for page_num, text in extracted_page_texts.items():
                displayed_page_num = printed_page_map.get(page_num) or page_labels.get(page_num) or page_num
                chunk_counter = self._append_chunks_from_text(
                    text=text,
                    doc_name=doc_name,
                    page_num=page_num,
                    chunk_counter=chunk_counter,
                    metadata={
                        "total_pages": len(pdf_reader.pages),
                        "processed_at": datetime.now().isoformat(),
                        "file_type": "pdf",
                        "source": "text",
                        "preview_kind": "pdf_page",
                        "preview_page_num": page_num,
                        "printed_page_num": displayed_page_num,
                        "preview_image": page_previews.get(page_num),
                    },
                )

            # OCR any embedded images in the PDF
            chunk_counter = self._extract_and_ocr_pdf_images(
                pdf_bytes,
                doc_name,
                chunk_counter,
                page_previews,
                printed_page_map,
                page_labels,
            )

        except PyPDF2.errors.PdfReadError as e:
            failed_files.append((doc_name, f"Invalid or corrupted PDF: {e}"))

        return chunk_counter

    def _extract_printed_page_candidate(self, text: str) -> Optional[int]:
        """Try to extract a human-visible printed page number from page footer text."""
        if not text:
            return None

        lines = [line.strip() for line in text.splitlines() if line and line.strip()]
        if not lines:
            return None

        # Footer numbers are usually in the last few non-empty lines.
        tail_lines = lines[-8:]
        patterns = [
            r"(?i)^chapter\s+\d+.*?\b(\d{1,4})\s*$",
            r"^(\d{1,4})$",
            r"\b(\d{1,4})\s*$",
        ]

        for line in reversed(tail_lines):
            lowered = line.lower()
            if "copyright" in lowered:
                continue

            for pattern in patterns:
                match = re.search(pattern, line)
                if not match:
                    continue

                try:
                    value = int(match.group(1))
                except (TypeError, ValueError):
                    continue

                if 1 <= value <= 5000:
                    return value

        return None

    def _extract_pdf_page_labels(self, pdf_bytes: bytes) -> Dict[int, int]:
        """Extract logical page labels from PDF metadata when available (e.g., label 35 on physical page 41)."""
        try:
            import fitz  # PyMuPDF
        except ImportError:
            return {}

        labels: Dict[int, int] = {}
        try:
            pdf_doc = fitz.open(stream=pdf_bytes, filetype="pdf")
            for page_num, page in enumerate(pdf_doc, start=1):
                raw_label = ""
                try:
                    if hasattr(page, "get_label"):
                        raw_label = page.get_label() or ""
                except Exception:
                    raw_label = ""

                if not raw_label:
                    continue

                match = re.search(r"(\d{1,4})$", str(raw_label).strip())
                if not match:
                    continue

                value = int(match.group(1))
                if value > 0:
                    labels[page_num] = value

            pdf_doc.close()
        except Exception as e:
            logger.debug(f"Failed to extract PDF page labels: {e}")

        return labels

    def _resolve_printed_page_map(self, candidates: Dict[int, int], total_pages: int) -> Dict[int, int]:
        """Resolve a stable PDF-page -> printed-page mapping from detected footer candidates."""
        if len(candidates) < 3:
            return {}

        offsets: Dict[int, int] = {}
        for pdf_page, printed_page in candidates.items():
            offset = pdf_page - printed_page
            offsets[offset] = offsets.get(offset, 0) + 1

        best_offset, support = max(offsets.items(), key=lambda item: item[1])
        if support < 3:
            return {}

        mapping: Dict[int, int] = {}
        for pdf_page in range(1, total_pages + 1):
            printed_page = pdf_page - best_offset
            if printed_page >= 1:
                mapping[pdf_page] = printed_page

        return mapping

    def _build_pdf_page_previews(self, pdf_bytes: bytes) -> Dict[int, bytes]:
        """Render small thumbnail previews for PDF pages."""
        try:
            import fitz  # PyMuPDF
            from PIL import Image
            import io
        except ImportError:
            logger.debug("PyMuPDF or Pillow not available – skipping PDF previews")
            return {}

        previews: Dict[int, bytes] = {}
        try:
            pdf_doc = fitz.open(stream=pdf_bytes, filetype="pdf")
            for page_num, page in enumerate(pdf_doc, start=1):
                try:
                    pixmap = page.get_pixmap(matrix=fitz.Matrix(1.6, 1.6), alpha=False)
                    image = Image.open(io.BytesIO(pixmap.tobytes("png"))).convert("RGB")
                    image.thumbnail((720, 960))

                    output = io.BytesIO()
                    image.save(output, format="JPEG", quality=78, optimize=True)
                    previews[page_num] = output.getvalue()
                except Exception as e:
                    logger.debug(f"Failed to render preview for PDF page {page_num}: {e}")
                    continue
            pdf_doc.close()
        except Exception as e:
            logger.warning(f"Failed to create PDF previews: {e}")

        return previews

    def _extract_and_ocr_pdf_images(
        self,
        pdf_bytes: bytes,
        doc_name: str,
        chunk_counter: int,
        page_previews: Optional[Dict[int, bytes]] = None,
        printed_page_map: Optional[Dict[int, int]] = None,
        page_labels: Optional[Dict[int, int]] = None,
    ) -> int:
        """
        Extract embedded images from a PDF using PyMuPDF and run Tesseract OCR.
        Silently skips if PyMuPDF or pytesseract are not installed.
        """
        try:
            import fitz  # PyMuPDF
            import pytesseract
            from PIL import Image
            import io
        except ImportError:
            logger.debug("PyMuPDF or pytesseract not available – skipping image OCR")
            return chunk_counter

        try:
            pdf_doc = fitz.open(stream=pdf_bytes, filetype="pdf")
            for page_num, page in enumerate(pdf_doc, start=1):
                for img in page.get_images(full=True):
                    try:
                        xref = img[0]
                        base_image = pdf_doc.extract_image(xref)
                        image_bytes = base_image["image"]

                        pil_img = Image.open(io.BytesIO(image_bytes)).convert("RGB")
                        ocr_text = pytesseract.image_to_string(pil_img).strip()

                        # Only keep OCR results with meaningful content
                        if len(ocr_text) < 50:
                            continue

                        chunk_counter = self._append_chunks_from_text(
                            text=ocr_text,
                            doc_name=doc_name,
                            page_num=page_num,
                            chunk_counter=chunk_counter,
                            metadata={
                                "processed_at": datetime.now().isoformat(),
                                "file_type": "pdf",
                                "source": "image_ocr",
                                "preview_kind": "pdf_page",
                                "preview_page_num": page_num,
                                "printed_page_num": (printed_page_map or {}).get(page_num)
                                or (page_labels or {}).get(page_num)
                                or page_num,
                                "preview_image": (page_previews or {}).get(page_num),
                            },
                        )
                    except Exception as e:
                        logger.debug(f"Failed to OCR image on page {page_num} of '{doc_name}': {e}")
                        continue
            pdf_doc.close()
        except Exception as e:
            logger.warning(f"Image OCR extraction failed for '{doc_name}': {e}")

        return chunk_counter

    def _process_docx_file(self, docx_file, doc_name: str, chunk_counter: int, failed_files: list) -> int:
        """Extract text and OCR images from a DOCX uploaded file and append chunks."""
        try:
            import io
            raw_bytes = docx_file.read()
            doc = docx.Document(io.BytesIO(raw_bytes))

            # Group paragraphs into virtual "pages" of ~2000 chars for page tracking
            full_text = "\n".join(para.text for para in doc.paragraphs if para.text.strip())

            if not full_text.strip():
                logger.warning(f"Word document '{doc_name}' has no text")
                failed_files.append((doc_name, "No text found"))
                return chunk_counter

            logger.info(f"Processing Word document '{doc_name}'")

            # Split into pseudo-pages of ~2000 chars
            page_size = 2000
            pages = [full_text[i:i + page_size] for i in range(0, len(full_text), page_size)]

            for page_num, page_text in enumerate(pages, start=1):
                chunk_counter = self._append_chunks_from_text(
                    text=page_text,
                    doc_name=doc_name,
                    page_num=page_num,
                    chunk_counter=chunk_counter,
                    metadata={
                        "total_pages": len(pages),
                        "processed_at": datetime.now().isoformat(),
                        "file_type": "docx",
                        "source": "text",
                    },
                )

            # OCR embedded images in the DOCX
            chunk_counter = self._extract_and_ocr_docx_images(doc, doc_name, chunk_counter)

        except Exception as e:
            failed_files.append((doc_name, f"Word document error: {e}"))

        return chunk_counter

    def _extract_and_ocr_docx_images(self, doc, doc_name: str, chunk_counter: int) -> int:
        """
        Extract embedded images from a python-docx Document and run Tesseract OCR.
        Silently skips if pytesseract / Pillow are not installed.
        """
        try:
            import pytesseract
            from PIL import Image
            import io
        except ImportError:
            logger.debug("pytesseract or Pillow not available – skipping DOCX image OCR")
            return chunk_counter

        try:
            for rel in doc.part.rels.values():
                if "image" not in rel.reltype:
                    continue
                try:
                    image_bytes = rel.target_part.blob
                    pil_img = Image.open(io.BytesIO(image_bytes)).convert("RGB")
                    ocr_text = pytesseract.image_to_string(pil_img).strip()

                    if len(ocr_text) < 50:
                        continue

                    chunk_counter = self._append_chunks_from_text(
                        text=ocr_text,
                        doc_name=doc_name,
                        page_num=1,
                        chunk_counter=chunk_counter,
                        metadata={
                            "processed_at": datetime.now().isoformat(),
                            "file_type": "docx",
                            "source": "image_ocr",
                        },
                    )
                except Exception as e:
                    logger.debug(f"Failed to OCR DOCX image in '{doc_name}': {e}")
                    continue
        except Exception as e:
            logger.warning(f"DOCX image OCR failed for '{doc_name}': {e}")

        return chunk_counter

    def _process_txt_file(self, txt_file, doc_name: str, chunk_counter: int, failed_files: list) -> int:
        """Extract text from a TXT uploaded file and append chunks."""
        try:
            raw_bytes = txt_file.read()
            decoded_text = None

            # Try common encodings in order.
            for enc in ("utf-8", "utf-16", "latin-1"):
                try:
                    decoded_text = raw_bytes.decode(enc)
                    break
                except UnicodeDecodeError:
                    continue

            if decoded_text is None or not decoded_text.strip():
                failed_files.append((doc_name, "No readable text found"))
                return chunk_counter

            logger.info(f"Processing text file '{doc_name}'")

            # Create pseudo-pages for parity with page-based document types.
            page_size = 3000
            pages = [decoded_text[i:i + page_size] for i in range(0, len(decoded_text), page_size)]

            for page_num, page_text in enumerate(pages, start=1):
                chunk_counter = self._append_chunks_from_text(
                    text=page_text,
                    doc_name=doc_name,
                    page_num=page_num,
                    chunk_counter=chunk_counter,
                    metadata={
                        "total_pages": len(pages),
                        "processed_at": datetime.now().isoformat(),
                        "file_type": "txt",
                        "source": "text",
                    },
                )

        except Exception as e:
            failed_files.append((doc_name, f"Text file error: {e}"))

        return chunk_counter

    def _process_pptx_file(self, pptx_file, doc_name: str, chunk_counter: int, failed_files: list) -> int:
        """Extract text from a PPTX uploaded file and append chunks."""
        try:
            import io
            from pptx import Presentation

            raw_bytes = pptx_file.read()
            presentation = Presentation(io.BytesIO(raw_bytes))

            if len(presentation.slides) == 0:
                failed_files.append((doc_name, "No slides found"))
                return chunk_counter

            logger.info(f"Processing PowerPoint '{doc_name}' ({len(presentation.slides)} slides)")

            for slide_num, slide in enumerate(presentation.slides, start=1):
                texts = []
                for shape in slide.shapes:
                    shape_text = getattr(shape, "text", None)
                    if shape_text and shape_text.strip():
                        texts.append(shape_text.strip())

                slide_text = "\n".join(texts).strip()
                chunk_counter = self._append_chunks_from_text(
                    text=slide_text,
                    doc_name=doc_name,
                    page_num=slide_num,
                    chunk_counter=chunk_counter,
                    metadata={
                        "total_pages": len(presentation.slides),
                        "processed_at": datetime.now().isoformat(),
                        "file_type": "pptx",
                        "source": "text",
                    },
                )

        except Exception as e:
            failed_files.append((doc_name, f"PowerPoint file error: {e}"))

        return chunk_counter

    
    def load_documents_from_directory(self, pdf_directory: str) -> int:
        """
        Load PDF documents from a directory.
        
        Args:
            pdf_directory: Path to directory containing PDFs
            
        Returns:
            Number of chunks created
            
        Raises:
            DocumentProcessingError: If directory doesn't exist or no valid PDFs found
        """
        pdf_path = Path(pdf_directory)
        
        if not pdf_path.exists():
            raise DocumentProcessingError(f"Directory does not exist: {pdf_directory}")
        
        if not pdf_path.is_dir():
            raise DocumentProcessingError(f"Path is not a directory: {pdf_directory}")
        
        pdf_files_paths = list(pdf_path.glob("*.pdf"))
        
        if not pdf_files_paths:
            raise DocumentProcessingError(f"No PDF files found in {pdf_directory}")
        
        logger.info(f"Found {len(pdf_files_paths)} PDF files in {pdf_directory}")
        
        self.documents.clear()
        self.embeddings_matrix = None
        self.index = None
        self.is_indexed = False
        chunk_counter = 0
        failed_files = []
        
        for pdf_file_path in tqdm(pdf_files_paths, desc="Processing PDFs"):
            try:
                with open(pdf_file_path, 'rb') as f:
                    pdf_reader = PyPDF2.PdfReader(f)
                    doc_name = pdf_file_path.name
                    
                    if len(pdf_reader.pages) == 0:
                        logger.warning(f"PDF '{doc_name}' has no pages")
                        failed_files.append((doc_name, "No pages found"))
                        continue
                    
                    for page_num, page in enumerate(pdf_reader.pages, start=1):
                        try:
                            text = page.extract_text()
                            chunk_counter = self._append_chunks_from_text(
                                text=text,
                                doc_name=doc_name,
                                page_num=page_num,
                                chunk_counter=chunk_counter,
                            )
                        
                        except Exception as e:
                            logger.error(f"Error processing page {page_num} of '{doc_name}': {e}")
                            continue
                            
            except (IOError, OSError) as e:
                error_msg = f"File access error: {e}"
                logger.error(f"Error reading {pdf_file_path}: {error_msg}")
                failed_files.append((pdf_file_path.name, error_msg))
                continue
                
            except PyPDF2.errors.PdfReadError as e:
                error_msg = f"Invalid PDF format: {e}"
                logger.error(f"Error processing {pdf_file_path}: {error_msg}")
                failed_files.append((pdf_file_path.name, error_msg))
                continue
                
            except Exception as e:
                error_msg = str(e)
                logger.error(f"Error processing {pdf_file_path}: {error_msg}")
                failed_files.append((pdf_file_path.name, error_msg))
                continue
        
        if failed_files:
            logger.warning(f"Failed to process {len(failed_files)} files")
        
        if len(self.documents) == 0:
            raise DocumentProcessingError(
                f"No valid documents could be processed. Failed files: {failed_files}"
            )
        
        logger.info(f"Created {len(self.documents)} chunks from {len(pdf_files_paths) - len(failed_files)} documents")
        return len(self.documents)
    
    def _create_chunks(self, text: str) -> List[str]:
        """
        Split text into overlapping chunks.
        
        Args:
            text: Text to chunk
            
        Returns:
            List of text chunks
        """
        # Clean text
        text = self._normalize_text_for_chunking(text)
        
        if len(text) <= self.chunk_size:
            return [text]
        
        chunks = []
        start = 0
        
        while start < len(text):
            end = start + self.chunk_size
            
            # Try to break at sentence boundary
            if end < len(text):
                # Look for sentence end
                last_period = text.rfind('.', start, end)
                if last_period > start + self.chunk_size // 2:
                    end = last_period + 1
            
            chunk = text[start:end].strip()
            if chunk:
                chunks.append(chunk)

            start = end - self.chunk_overlap
            # Advance start to the next word boundary so chunks never begin mid-word
            while start < len(text) and start > 0 and text[start - 1] not in (' ', '\n'):
                start += 1
        
        return chunks

    def _normalize_text_for_chunking(self, text: str) -> str:
        """Normalize OCR/book text to improve chunk quality before indexing."""
        normalized = text or ""
        normalized = normalized.replace("\r", "\n")

        # Fix merged boundaries like "AIAmong" that appear in extracted PDF text.
        normalized = re.sub(r"([A-Z]{2,})([A-Z][a-z])", r"\1 \2", normalized)

        # Repair split words from PDF extraction/OCR.
        normalized = re.sub(r"(\w)-\s+(\w)", r"\1\2", normalized)

        # Remove URLs and common book-preview/footer clutter.
        normalized = re.sub(r"https?://\S+", " ", normalized, flags=re.IGNORECASE)
        normalized = re.sub(r"\bMore books\s*:?[^\n]*", " ", normalized, flags=re.IGNORECASE)

        # Remove leading all-caps title/header fragments (generic, not book-specific).
        normalized = re.sub(
            r"^(?:\s*[A-Z][A-Z0-9&'()/:;,-]{2,}(?:\s+[A-Z][A-Z0-9&'()/:;,-]{2,}){3,18})(?=\s+[A-Z][a-z]|\s+[A-Z]{2,}\b)",
            " ",
            normalized,
        )

        # Collapse excessive whitespace.
        normalized = re.sub(r"\s+", " ", normalized).strip()
        return normalized

    def _query_terms(self, query: str) -> set:
        """Extract meaningful lexical terms from query for lightweight reranking."""
        tokens = re.findall(r"[A-Za-z0-9]+", query.lower())
        return {token for token in tokens if len(token) >= 3 and token not in _QUERY_STOPWORDS}

    def _lexical_overlap_score(self, query_terms: set, text: str) -> float:
        """Compute lexical overlap in range [0, 1] between query terms and chunk text."""
        if not query_terms:
            return 0.0
        words = set(re.findall(r"[A-Za-z0-9]+", text.lower()))
        if not words:
            return 0.0
        overlap = len(query_terms & words)
        return overlap / len(query_terms)

    def _noise_penalty(self, text: str) -> float:
        """Penalize chunk text that looks like headings/authors instead of explanatory content."""
        penalty = 0.0
        collapsed = re.sub(r"\s+", " ", text).strip()
        if not collapsed:
            return 0.2

        words = re.findall(r"[A-Za-z']+", collapsed)
        if not words:
            return 0.2

        upper_ratio = sum(1 for word in words if word.isupper() and len(word) > 2) / len(words)
        if upper_ratio > 0.35:
            penalty += 0.08

        lowered = collapsed.lower()
        if "use case" in lowered and "retrieval augmented generation" in lowered:
            penalty += 0.03
        if lowered.startswith("chapter") or lowered.startswith("table of contents"):
            penalty += 0.05

        return min(penalty, 0.2)

    def _low_information_penalty(self, text: str) -> float:
        """Penalize fragment-like snippets that are too short or incomplete to be useful evidence."""
        collapsed = re.sub(r"\s+", " ", (text or "")).strip()
        if not collapsed:
            return 0.25

        words = re.findall(r"[A-Za-z0-9']+", collapsed)
        if len(collapsed) < _MIN_RESULT_CHARS or len(words) < _MIN_RESULT_WORDS:
            return 0.25

        # Mid-word tails like "... which ha" are strong indicators of broken chunks.
        if len(words) >= 1 and collapsed[-1].isalpha() and collapsed[-1] not in ".!?":
            if len(collapsed) < 120:
                return 0.15

        return 0.0
    
    def generate_embeddings(self, batch_size: int = 32, show_progress: bool = True) -> None:
        """
        Generate embeddings for all document chunks.
        
        Args:
            batch_size: Batch size for embedding generation
            show_progress: Whether to show progress bar
            
        Raises:
            EmbeddingGenerationError: If embedding generation fails
            ValueError: If no documents loaded
        """
        if not self.documents:
            raise ValueError("No documents loaded. Call load_documents first.")
        
        if batch_size <= 0:
            raise ValueError("batch_size must be positive")
        
        logger.info(f"Generating embeddings for {len(self.documents)} chunks")
        
        try:
            # Extract texts
            texts = [doc.text for doc in self.documents]
            
            # Generate embeddings in batches
            all_embeddings = self.embedding_model.encode(
                texts,
                batch_size=batch_size,
                show_progress_bar=show_progress,
                convert_to_numpy=True,
                normalize_embeddings=True  # For cosine similarity
            )
            
            # Validate embeddings
            if all_embeddings is None or len(all_embeddings) == 0:
                raise EmbeddingGenerationError("Embedding generation returned empty results")
            
            if len(all_embeddings) != len(self.documents):
                raise EmbeddingGenerationError(
                    f"Embedding count mismatch: expected {len(self.documents)}, got {len(all_embeddings)}"
                )

            self.embeddings_matrix = all_embeddings.astype("float32", copy=False)
            if self.embeddings_matrix.shape[1] != self.embedding_dim:
                raise EmbeddingGenerationError(
                    f"Embedding dimension mismatch: expected {self.embedding_dim}, got {self.embeddings_matrix.shape[1]}"
                )

            # Keep chunk metadata lightweight once matrix is prepared.
            for doc in self.documents:
                doc.embedding = None
            
            logger.info("Embeddings generated successfully")
            
        except Exception as e:
            logger.error(f"Failed to generate embeddings: {e}")
            raise EmbeddingGenerationError(f"Embedding generation failed: {e}") from e
    
    def build_index(self) -> None:
        """
        Build FAISS index for efficient similarity search.
        
        Raises:
            IndexBuildError: If index building fails
            ValueError: If embeddings not generated
        """
        if not self.documents:
            raise ValueError("No documents loaded. Call load_documents first.")
        
        if self.embeddings_matrix is None or len(self.embeddings_matrix) == 0:
            raise ValueError("Generate embeddings before building index")
        
        logger.info("Building FAISS index")
        
        try:
            # Create embedding matrix
            embeddings_matrix = self.embeddings_matrix
            
            # Validate embedding matrix
            if embeddings_matrix.shape[0] != len(self.documents):
                raise IndexBuildError(f"Embedding matrix shape mismatch")
            
            if embeddings_matrix.shape[1] != self.embedding_dim:
                raise IndexBuildError(
                    f"Embedding dimension mismatch: expected {self.embedding_dim}, got {embeddings_matrix.shape[1]}"
                )
            
            # Build index (Inner Product = Cosine similarity for normalized vectors)
            self.index = faiss.IndexFlatIP(self.embedding_dim)
            self.index.add(embeddings_matrix)
            self.is_indexed = True
            
            logger.info(f"Index built with {self.index.ntotal} vectors")
            
        except Exception as e:
            logger.error(f"Failed to build FAISS index: {e}")
            self.is_indexed = False
            raise IndexBuildError(f"Index building failed: {e}") from e
    
    def retrieve(
        self,
        query: str,
        top_k: int = 5,
        min_score: Optional[float] = None,
    ) -> List[Tuple[DocumentChunk, float]]:
        """
        Retrieve most relevant documents for a query.
        
        Args:
            query: User query string
            top_k: Number of documents to retrieve
            
        Returns:
            List of (DocumentChunk, similarity_score) tuples
            
        Raises:
            SearchError: If search fails
            ValueError: If invalid parameters or index not built
        """
        if not self.is_indexed:
            raise ValueError("Index not built. Call build_index first.")
        
        if not query or not query.strip():
            raise ValueError("Query cannot be empty")
        
        if top_k <= 0:
            raise ValueError("top_k must be positive")

        if min_score is not None and not 0 <= min_score <= 1:
            raise ValueError("min_score must be between 0 and 1")
        
        # Limit top_k to available documents
        top_k = min(top_k, len(self.documents))
        
        try:
            # Generate query embedding
            query_embedding = self.embedding_model.encode(
                [query],
                convert_to_numpy=True,
                normalize_embeddings=True
            ).astype('float32')
            
            # Validate query embedding
            if query_embedding is None or len(query_embedding) == 0:
                raise SearchError("Failed to generate query embedding")
            
            # Search extra candidates then rerank with lexical overlap and noise penalty.
            search_k = min(len(self.documents), max(top_k * 4, top_k + 8))
            scores, indices = self.index.search(query_embedding, search_k)

            query_terms = self._query_terms(query)
            reranked = []
            for idx, vector_score in zip(indices[0], scores[0]):
                if idx < 0:
                    continue
                chunk = self.documents[idx]
                low_info_penalty = self._low_information_penalty(chunk.text)
                if low_info_penalty >= 0.25:
                    continue

                lexical_score = self._lexical_overlap_score(query_terms, chunk.text)

                # Guardrail: if lexical overlap is absent and semantic score is not very strong,
                # treat it as likely off-topic for user-facing ranking.
                if query_terms and lexical_score == 0.0 and float(vector_score) < 0.70:
                    continue

                noise_penalty = self._noise_penalty(chunk.text)
                combined_score = (0.72 * float(vector_score)) + (0.38 * lexical_score) - noise_penalty - low_info_penalty
                combined_score = max(0.0, min(1.0, combined_score))
                reranked.append((chunk, combined_score))

            results = sorted(reranked, key=lambda item: item[1], reverse=True)[:top_k]

            if min_score is not None:
                results = [(doc, score) for doc, score in results if score >= min_score]
            
            logger.info(f"Retrieved {len(results)} documents for query")
            return results
            
        except Exception as e:
            logger.error(f"Search failed: {e}")
            raise SearchError(f"Failed to retrieve documents: {e}") from e
    
    def format_context_with_citations(self, retrieved_docs: List[Tuple[DocumentChunk, float]]) -> Tuple[str, List[Dict]]:
        """
        Format retrieved documents into context with citation markers.
        
        Args:
            retrieved_docs: List of (DocumentChunk, score) tuples
            
        Returns:
            Tuple of (formatted_context, citations_list)
        """
        context_parts = []
        citations = []
        
        for idx, (doc, score) in enumerate(retrieved_docs, start=1):
            citation_marker = f"[{idx}]"
            context_parts.append(f"{citation_marker} {doc.text}")

            citations.append({
                "id": idx,
                "document": doc.doc_name,
                "page": doc.page_num,
                "score": score,
                "source": doc.metadata.get("source", "text"),
                "text_preview": doc.text[:200] + "..." if len(doc.text) > 200 else doc.text
            })
        
        context = "\n\n".join(context_parts)
        return context, citations
    
    def get_statistics(self) -> Dict[str, Any]:
        """Get pipeline statistics."""
        if not self.documents:
            return {"status": "No documents loaded"}
        
        doc_names = [doc.doc_name for doc in self.documents]
        unique_docs = set(doc_names)
        
        return {
            "total_chunks": len(self.documents),
            "total_documents": len(unique_docs),
            "documents": list(unique_docs),
            "chunks_per_doc": {doc: doc_names.count(doc) for doc in unique_docs},
            "embedding_dim": self.embedding_dim,
            "is_indexed": self.is_indexed,
            "chunk_size": self.chunk_size,
            "chunk_overlap": self.chunk_overlap
        }
