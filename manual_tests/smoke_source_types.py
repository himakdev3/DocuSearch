from __future__ import annotations

from dataclasses import dataclass
from pathlib import Path

import fitz
from docx import Document
from pptx import Presentation

from src.retrieval.rag_pipeline import RAGPipeline


@dataclass
class MockUpload:
    name: str
    _data: bytes

    def read(self) -> bytes:
        return self._data


def _make_txt(path: Path) -> None:
    path.write_text(
        "TXT source content: oak taxonomy and deciduous species overview.",
        encoding="utf-8",
    )


def _make_docx(path: Path) -> None:
    doc = Document()
    doc.add_heading("DOCX source", level=1)
    doc.add_paragraph("DOCX source content: river ecology and riparian habitats.")
    doc.save(path)


def _make_pptx(path: Path) -> None:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "PPTX source"
    slide.placeholders[1].text = "PPTX source content: climate adaptation and resilience."
    prs.save(path)


def _make_pdf(path: Path) -> None:
    doc = fitz.open()
    page = doc.new_page()
    page.insert_text((72, 100), "PDF source content: conifer classification and needle traits.")
    doc.save(path)
    doc.close()


def _read_upload(path: Path) -> MockUpload:
    return MockUpload(name=path.name, _data=path.read_bytes())


def main() -> None:
    base = Path("manual_tests/generated_sources")
    base.mkdir(parents=True, exist_ok=True)

    txt = base / "sample_trees.txt"
    docx = base / "sample_trees.docx"
    pptx = base / "sample_trees.pptx"
    pdf = base / "sample_trees.pdf"

    _make_txt(txt)
    _make_docx(docx)
    _make_pptx(pptx)
    _make_pdf(pdf)

    uploads = [_read_upload(p) for p in [txt, docx, pptx, pdf]]

    pipeline = RAGPipeline(chunk_size=400, chunk_overlap=40)
    chunks = pipeline.load_documents_from_files(uploads)
    pipeline.generate_embeddings(batch_size=8, show_progress=False)
    pipeline.build_index()

    print(f"Created chunks: {chunks}")
    print("Chunks per doc:", pipeline.get_statistics().get("chunks_per_doc", {}))

    queries = {
        "deciduous species overview": "sample_trees.txt",
        "riparian habitats": "sample_trees.docx",
        "climate adaptation resilience": "sample_trees.pptx",
        "conifer needle traits": "sample_trees.pdf",
    }

    failures = []
    for query, expected_doc in queries.items():
        results = pipeline.retrieve(query, top_k=4, min_score=0.2)
        docs = [doc.doc_name for doc, _ in results]
        print(f"Query: {query}")
        print("Top docs:", docs)
        if expected_doc not in docs:
            failures.append((query, expected_doc, docs))

    if failures:
        print("\nSMOKE TEST FAILED")
        for query, expected_doc, docs in failures:
            print(f"- Missing {expected_doc} for query '{query}'. Got: {docs}")
        raise SystemExit(1)

    print("\nSMOKE TEST PASSED: all source types were indexed and retrievable.")


if __name__ == "__main__":
    main()
