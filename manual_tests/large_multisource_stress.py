from __future__ import annotations

from dataclasses import dataclass
from pathlib import Path
from time import perf_counter

import fitz
from docx import Document
from pptx import Presentation

from src.retrieval.rag_pipeline import RAGPipeline


DOCX_PAGE_SIZE = 2000
TXT_PAGE_SIZE = 3000


@dataclass
class MockUpload:
    name: str
    _data: bytes

    def read(self) -> bytes:
        return self._data


def _paragraph(seed: str, repeat: int = 8) -> str:
    unit = (
        f"{seed} This section describes tree form, root behavior, canopy exchange, soil moisture, "
        "species interactions, and long-horizon restoration planning with citations and examples. "
    )
    return unit * repeat


def _build_large_txt(path: Path, sections: int = 14) -> None:
    blocks = []
    for i in range(1, sections + 1):
        marker = f"TXT_MARKER_SECTION_{i:02d}"
        blocks.append(_paragraph(f"Text section {i}. {marker}.", repeat=20))
    path.write_text("\n\n".join(blocks), encoding="utf-8")


def _build_large_docx(path: Path, sections: int = 14) -> None:
    doc = Document()
    doc.add_heading("Large DOCX Multi-Page Source", level=1)
    for i in range(1, sections + 1):
        marker = f"DOCX_MARKER_SECTION_{i:02d}"
        doc.add_heading(f"Section {i}", level=2)
        doc.add_paragraph(_paragraph(f"Docx section {i}. {marker}.", repeat=16))
        doc.add_page_break()
    doc.save(path)


def _build_large_pptx(path: Path, slides: int = 18) -> None:
    prs = Presentation()
    for i in range(1, slides + 1):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        marker = f"PPTX_MARKER_SLIDE_{i:02d}"
        slide.shapes.title.text = f"Slide {i}: Forest Analytics"
        slide.placeholders[1].text = _paragraph(
            f"Presentation slide {i}. {marker}.", repeat=10
        )
    prs.save(path)


def _build_large_pdf(path: Path, pages: int = 16) -> None:
    doc = fitz.open()
    for i in range(1, pages + 1):
        page = doc.new_page()
        marker = f"PDF_MARKER_PAGE_{i:02d}"
        text = _paragraph(f"PDF page {i}. {marker}.", repeat=14)
        page.insert_textbox(fitz.Rect(50, 60, 545, 780), text, fontsize=11)
    doc.save(path)
    doc.close()


def _to_upload(path: Path) -> MockUpload:
    return MockUpload(path.name, path.read_bytes())


def main() -> None:
    root = Path("manual_tests/generated_large_sources")
    root.mkdir(parents=True, exist_ok=True)

    txt = root / "large_sample.txt"
    docx = root / "large_sample.docx"
    pptx = root / "large_sample.pptx"
    pdf = root / "large_sample.pdf"

    print("Generating large multi-section sources...")
    _build_large_txt(txt)
    _build_large_docx(docx)
    _build_large_pptx(pptx)
    _build_large_pdf(pdf)

    files = [txt, docx, pptx, pdf]
    for f in files:
        print(f"- {f.name}: {f.stat().st_size / 1024:.1f} KB")

    uploads = [_to_upload(f) for f in files]

    pipeline = RAGPipeline(chunk_size=500, chunk_overlap=50)

    t0 = perf_counter()
    chunk_count = pipeline.load_documents_from_files(uploads)
    t1 = perf_counter()

    pipeline.generate_embeddings(batch_size=16, show_progress=False)
    t2 = perf_counter()

    pipeline.build_index()
    t3 = perf_counter()

    stats = pipeline.get_statistics()
    print("\nIndex statistics:")
    print("- total chunks:", chunk_count)
    print("- chunks per doc:", stats.get("chunks_per_doc", {}))
    print("- load time: %.2fs" % (t1 - t0))
    print("- embedding time: %.2fs" % (t2 - t1))
    print("- index time: %.2fs" % (t3 - t2))

    # Validate retrieval from deep sections/pages/slides.
    checks = {
        "TXT_MARKER_SECTION_13": "large_sample.txt",
        "DOCX_MARKER_SECTION_12": "large_sample.docx",
        "PPTX_MARKER_SLIDE_17": "large_sample.pptx",
        "PDF_MARKER_PAGE_15": "large_sample.pdf",
    }

    failures = []
    for query, expected in checks.items():
        res = pipeline.retrieve(query, top_k=6, min_score=0.2)
        docs = [doc.doc_name for doc, _ in res]
        tops = [(doc.doc_name, doc.page_num, round(score, 3)) for doc, score in res[:3]]
        print(f"\nQuery: {query}")
        print("Top 3:", tops)
        if expected not in docs:
            failures.append((query, expected, docs))

    if failures:
        print("\nLARGE MULTI-SOURCE TEST FAILED")
        for q, exp, docs in failures:
            print(f"- Expected {exp} for '{q}' but got: {docs}")
        raise SystemExit(1)

    print("\nLARGE MULTI-SOURCE TEST PASSED")


if __name__ == "__main__":
    main()
